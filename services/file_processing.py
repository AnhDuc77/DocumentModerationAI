"""
File processing services for AI Document Processing Service
"""

import io
import base64
from typing import Dict, Any, List, Tuple
from PIL import Image
import PyPDF2
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from io import BytesIO

# PDF image extraction libraries
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("PyMuPDF not available. Install with: pip install PyMuPDF")

try:
    from pdf2image import convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    print("pdf2image not available. Install with: pip install pdf2image")

from services.text_moderation import moderate_educational_content
from services.image_moderation import moderate_image_advanced
from utils.file_utils import detect_file_type, is_supported_file_type, validate_file_size
from config.settings import FILE_PROCESSING_THRESHOLDS

def split_text_into_chunks(text: str, chunk_size: int = 2000) -> List[str]:
    """Smart text chunking: only chunk large files for optimal performance"""
    if not text or len(text) <= chunk_size:
        return [text] if text else []
    
    chunks = []
    words = text.split()
    current_chunk = []
    current_length = 0
    
    for word in words:
        word_length = len(word) + 1  # +1 for space
        
        # If adding this word would exceed chunk size, start new chunk
        if current_length + word_length > chunk_size and current_chunk:
            chunks.append(' '.join(current_chunk))
            current_chunk = [word]
            current_length = len(word)
        else:
            current_chunk.append(word)
            current_length += word_length
    
    # Add the last chunk if it has content
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks

def smart_text_chunking(text: str) -> List[str]:
    """Smart chunking with dilution protection for optimal performance and accuracy"""
    # Small files: process as single chunk (faster)
    if len(text) <= 1500:
        return [text] if text else []
    
    # Medium files: smaller chunks to avoid dilution (better detection)
    elif len(text) <= 8000:
        return split_text_into_chunks(text, chunk_size=1500)
    
    # Large files: standard chunks with dilution protection
    else:
        return split_text_into_chunks(text, chunk_size=1500)

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file"""
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        print(f"PDF text extraction failed: {e}")
        return ""

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from DOCX file"""
    try:
        doc = Document(io.BytesIO(file_content))
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"DOCX text extraction failed: {e}")
        return ""

def extract_text_from_xlsx(file_content: bytes) -> str:
    """Extract text from XLSX file"""
    try:
        workbook = load_workbook(io.BytesIO(file_content))
        text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        text += str(cell) + " "
                text += "\n"
        return text.strip()
    except Exception as e:
        print(f"XLSX text extraction failed: {e}")
        return ""

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PPTX file"""
    try:
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"PPTX text extraction failed: {e}")
        return ""

def extract_images_from_pdf(file_content: bytes) -> List[str]:
    """Extract images from PDF file using PyMuPDF and pdf2image"""
    images = []
    
    try:
        # Method 1: Extract embedded images using PyMuPDF (preferred)
        if PYMUPDF_AVAILABLE:
            images.extend(_extract_embedded_images_pymupdf(file_content))
        
        # Method 2: Convert PDF pages to images using pdf2image (fallback)
        if not images and PDF2IMAGE_AVAILABLE:
            images.extend(_extract_pages_as_images_pdf2image(file_content))
        
        return images
        
    except Exception as e:
        print(f"PDF image extraction failed: {e}")
        return []

def _extract_embedded_images_pymupdf(file_content: bytes) -> List[str]:
    """Extract embedded images from PDF using PyMuPDF"""
    images = []
    
    try:
        # Open PDF document
        doc = fitz.open(stream=file_content, filetype="pdf")
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Get list of images on this page
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                try:
                    # Get image data
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    # Skip if image is too small (likely decorative)
                    if pix.width < 50 or pix.height < 50:
                        pix = None
                        continue
                    
                    # Convert to PIL Image
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_data = pix.tobytes("png")
                        pil_image = Image.open(BytesIO(img_data))
                        
                        # Convert to base64
                        img_buffer = BytesIO()
                        pil_image.save(img_buffer, format='PNG')
                        img_base64 = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
                        images.append(img_base64)
                    
                    pix = None  # Free memory
                    
                except Exception as e:
                    print(f"Error extracting image {img_index} from page {page_num}: {e}")
                    continue
        
        doc.close()
        return images
        
    except Exception as e:
        print(f"PyMuPDF image extraction failed: {e}")
        return []

def _extract_pages_as_images_pdf2image(file_content: bytes) -> List[str]:
    """Convert PDF pages to images using pdf2image"""
    images = []
    
    try:
        # Convert PDF to images (one per page)
        pdf_images = convert_from_bytes(
            file_content,
            dpi=150,  # Good balance between quality and size
            fmt='PNG',
            first_page=1,
            last_page=None  # All pages
        )
        
        for page_num, pil_image in enumerate(pdf_images):
            try:
                # Skip if image is too small
                if pil_image.width < 100 or pil_image.height < 100:
                    continue
                
                # Convert to base64
                img_buffer = BytesIO()
                pil_image.save(img_buffer, format='PNG')
                img_base64 = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
                images.append(img_base64)
                
            except Exception as e:
                print(f"Error processing page {page_num + 1}: {e}")
                continue
        
        return images
        
    except Exception as e:
        print(f"pdf2image conversion failed: {e}")
        return []

def extract_images_from_docx(file_content: bytes) -> List[str]:
    """Extract images from DOCX file"""
    try:
        doc = Document(io.BytesIO(file_content))
        images = []
        
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                image_base64 = base64.b64encode(image_data).decode('utf-8')
                images.append(image_base64)
        
        return images
    except Exception as e:
        print(f"DOCX image extraction failed: {e}")
        return []

def extract_file_content(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Extract text and images from various file types"""
    file_type = detect_file_type(file_content, filename)
    
    result = {
        "file_type": file_type,
        "text_content": "",
        "images": [],
        "extraction_status": "success",
        "error": None
    }
    
    try:
        if not is_supported_file_type(filename):
            result["extraction_status"] = "unsupported_file_type"
            result["error"] = f"Unsupported file type: {file_type}"
            return result
        
        if not validate_file_size(len(file_content)):
            result["extraction_status"] = "file_too_large"
            result["error"] = "File size exceeds maximum allowed size"
            return result
        
        # Extract text based on file type
        if file_type == "pdf":
            result["text_content"] = extract_text_from_pdf(file_content)
            result["images"] = extract_images_from_pdf(file_content)
            result["extraction_details"] = {
                "pymupdf_available": PYMUPDF_AVAILABLE,
                "pdf2image_available": PDF2IMAGE_AVAILABLE,
                "image_extraction_method": "PyMuPDF" if PYMUPDF_AVAILABLE else ("pdf2image" if PDF2IMAGE_AVAILABLE else "none")
            }
        elif file_type == "docx":
            result["text_content"] = extract_text_from_docx(file_content)
            result["images"] = extract_images_from_docx(file_content)
        elif file_type in ["xlsx", "xls"]:
            result["text_content"] = extract_text_from_xlsx(file_content)
        elif file_type == "pptx":
            result["text_content"] = extract_text_from_pptx(file_content)
        else:
            result["extraction_status"] = "unsupported_file_type"
            result["error"] = f"Unsupported file type: {file_type}"
        
        return result
        
    except Exception as e:
        result["extraction_status"] = "extraction_error"
        result["error"] = str(e)
        return result

def moderate_file_content(text: str, images: List[str], thresholds: Dict[str, float] = None, 
                         include_detailed_results: bool = True) -> Dict[str, Any]:
    """Moderate extracted file content with chunk-based processing and page tracking"""
    if thresholds is None:
        thresholds = FILE_PROCESSING_THRESHOLDS.copy()
    
    result = {
        "text_moderation": {},
        "text_chunks": [],
        "image_moderation": [],
        "overall_score": 0.0,
        "overall_risk": "LOW",
        "thresholds_used": thresholds,
        "processing_details": {},
        "violation_summary": [],
        "violation_pages": [],
        "total_pages": 0,
        "violations_found": 0
    }
    
    try:
        # Process text content with smart chunking for optimal performance
        if text and text.strip():
            # Use smart chunking: only chunk when necessary
            text_chunks = smart_text_chunking(text)
            result["processing_details"]["text_chunks_count"] = len(text_chunks)
            result["processing_details"]["chunking_strategy"] = "smart"
            
            chunk_results = []
            all_text_scores = []
            
            # Process each chunk separately
            for i, chunk in enumerate(text_chunks):
                if chunk.strip():  # Skip empty chunks
                    chunk_result = moderate_educational_content(chunk)
                    chunk_result["chunk_index"] = i
                    chunk_result["chunk_text"] = chunk[:100] + "..." if len(chunk) > 100 else chunk
                    
                    # Estimate page number (rough calculation: 800 chars per page)
                    estimated_page = (i * 800) // 800 + 1
                    chunk_result["estimated_page"] = estimated_page
                    
                    chunk_results.append(chunk_result)
                    
                    # Check for violations in this chunk
                    max_chunk_score = max([
                        chunk_result.get("inappropriate_score", 0.0),
                        chunk_result.get("offensive_score", 0.0),
                        chunk_result.get("sexual_content_score", 0.0),
                        chunk_result.get("violence_score", 0.0)
                    ])
                    
                    if max_chunk_score > thresholds.get("overall", 0.8) * 0.5:  # 50% of threshold
                        violation_info = {
                            "page_number": estimated_page,
                            "violation_type": "text",
                            "severity": "HIGH" if max_chunk_score > thresholds.get("overall", 0.8) else "MEDIUM",
                            "score": max_chunk_score,
                            "details": f"Text violation in chunk {i+1}"
                        }
                        result["violation_summary"].append(violation_info)
                        if estimated_page not in result["violation_pages"]:
                            result["violation_pages"].append(estimated_page)
                    
                    # Collect scores from this chunk
                    chunk_scores = [
                        chunk_result.get("inappropriate_score", 0.0),
                        chunk_result.get("offensive_score", 0.0),
                        chunk_result.get("sexual_content_score", 0.0),
                        chunk_result.get("violence_score", 0.0)
                    ]
                    all_text_scores.extend(chunk_scores)
            
            result["text_chunks"] = chunk_results
            
            # Get the highest scores from all chunks (anti-dilution)
            if all_text_scores:
                max_text_score = max(all_text_scores)
                result["text_moderation"] = {
                    "inappropriate_score": max([r.get("inappropriate_score", 0.0) for r in chunk_results]),
                    "offensive_score": max([r.get("offensive_score", 0.0) for r in chunk_results]),
                    "sexual_content_score": max([r.get("sexual_content_score", 0.0) for r in chunk_results]),
                    "violence_score": max([r.get("violence_score", 0.0) for r in chunk_results]),
                    "overall_risk": "HIGH" if max_text_score > thresholds.get("overall", 0.8) else 
                                  "MEDIUM" if max_text_score > thresholds.get("overall", 0.8) * 0.7 else "LOW",
                    "chunks_processed": len(chunk_results),
                    "max_chunk_score": max_text_score
                }
                
                if max_text_score > thresholds.get("overall", 0.8):
                    result["overall_risk"] = "HIGH"
                elif max_text_score > thresholds.get("overall", 0.8) * 0.7:
                    result["overall_risk"] = "MEDIUM"
        
        # Moderate images
        if images:
            image_scores = []
            for i, image_base64 in enumerate(images):
                try:
                    # Decode base64 image
                    image_data = base64.b64decode(image_base64)
                    image = Image.open(BytesIO(image_data))
                    
                    # Moderate image
                    image_result = moderate_image_advanced(image)
                    image_result["image_index"] = i
                    
                    # Estimate page number for image (assuming images are in order)
                    estimated_page = i + 1
                    image_result["estimated_page"] = estimated_page
                    
                    result["image_moderation"].append(image_result)
                    
                    # Check for violations in this image
                    max_image_score = max([
                        image_result.get("nsfw_score", 0.0),
                        image_result.get("violence_score", 0.0)
                    ])
                    
                    if max_image_score > thresholds.get("overall", 0.8) * 0.5:  # 50% of threshold
                        violation_info = {
                            "page_number": estimated_page,
                            "violation_type": "image",
                            "severity": "HIGH" if max_image_score > thresholds.get("overall", 0.8) else "MEDIUM",
                            "score": max_image_score,
                            "details": f"Image violation in image {i+1}"
                        }
                        result["violation_summary"].append(violation_info)
                        if estimated_page not in result["violation_pages"]:
                            result["violation_pages"].append(estimated_page)
                    
                    # Collect scores
                    image_scores.extend([
                        image_result.get("nsfw_score", 0.0),
                        image_result.get("violence_score", 0.0)
                    ])
                    
                except Exception as e:
                    print(f"Image {i} moderation failed: {e}")
                    result["image_moderation"].append({
                        "image_index": i,
                        "error": str(e),
                        "nsfw_score": 0.0,
                        "violence_score": 0.0,
                        "estimated_page": i + 1
                    })
            
            # Check image scores against thresholds
            if image_scores:
                max_image_score = max(image_scores)
                if max_image_score > thresholds.get("overall", 0.8):
                    result["overall_risk"] = "HIGH"
                elif max_image_score > thresholds.get("overall", 0.8) * 0.7:
                    if result["overall_risk"] == "LOW":
                        result["overall_risk"] = "MEDIUM"
        
        # Calculate overall score from chunks (anti-dilution)
        all_scores = []
        
        # Add text scores from chunks
        if result["text_chunks"]:
            for chunk_result in result["text_chunks"]:
                all_scores.extend([
                    chunk_result.get("inappropriate_score", 0.0),
                    chunk_result.get("offensive_score", 0.0),
                    chunk_result.get("sexual_content_score", 0.0),
                    chunk_result.get("violence_score", 0.0)
                ])
        
        # Add image scores
        for img_result in result["image_moderation"]:
            all_scores.extend([
                img_result.get("nsfw_score", 0.0),
                img_result.get("violence_score", 0.0)
            ])
        
        result["overall_score"] = max(all_scores) if all_scores else 0.0
        
        # Calculate summary statistics
        result["total_pages"] = max(
            len(text_chunks) if text_chunks else 0,
            len(images) if images else 0,
            1  # At least 1 page
        )
        result["violations_found"] = len(result["violation_pages"])
        result["violation_pages"].sort()
        
        # Final risk assessment
        if result["overall_score"] > thresholds.get("overall", 0.8):
            result["overall_risk"] = "HIGH"
        elif result["overall_score"] > thresholds.get("overall", 0.8) * 0.6:
            result["overall_risk"] = "MEDIUM"
        else:
            result["overall_risk"] = "LOW"
        
        # Limit detailed results if requested
        if not include_detailed_results:
            result["text_chunks"] = result["text_chunks"][:5]  # Only first 5 chunks
            result["image_moderation"] = result["image_moderation"][:5]  # Only first 5 images
            result["text_content"] = result.get("text_content", "")[:500] + "..." if len(result.get("text_content", "")) > 500 else result.get("text_content", "")
        
        return result
        
    except Exception as e:
        result["overall_risk"] = "ERROR"
        result["processing_details"]["error"] = str(e)
        return result
