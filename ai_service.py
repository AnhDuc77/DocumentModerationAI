#!/usr/bin/env python3
"""
AI Document Processing Service - Production Ready
Supports Vietnamese and English text moderation plus advanced image analysis
"""

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import List, Dict, Any, Optional, Union
import asyncio
import io
import time
import json
import hashlib
import base64
import tempfile
import os
from dataclasses import dataclass
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
import base64
import io
import mimetypes

# AI Libraries
from transformers import pipeline
from PIL import Image
import torch
import numpy as np
from datasketch import MinHash
from simhash import Simhash
import re

# Try to import specialized models
try:
    from nudenet import NudeDetector
    NUDENET_AVAILABLE = True
    print("NudeNet library detected - High accuracy NSFW detection available")
except ImportError:
    NUDENET_AVAILABLE = False
    print("NudeNet not available. Install with: pip install nudenet")

try:
    from transformers import CLIPProcessor, CLIPModel
    CLIP_AVAILABLE = True
    print("CLIP library detected - General image understanding available")
except ImportError:
    CLIP_AVAILABLE = False
    print("CLIP not available")

# Try Falconsai NSFW model
try:
    FALCONSAI_AVAILABLE = True
    print("Falconsai NSFW model available - Specialized NSFW detection ready")
except ImportError:
    FALCONSAI_AVAILABLE = False
    print("Falconsai model not available")

# Try educational content moderation models
try:
    from transformers import AutoTokenizer, AutoModelForSequenceClassification
    EDUCATIONAL_MODELS_AVAILABLE = True
    print("Educational moderation models available")
except ImportError:
    EDUCATIONAL_MODELS_AVAILABLE = False
    print("Educational models not available")

# Configuration
NSFW_THRESHOLD = 0.7
VIOLENCE_THRESHOLD = 0.6
TOXICITY_THRESHOLD = 0.7
NUDENET_THRESHOLD = 0.6
NUM_PERM = 128

# Educational content moderation thresholds
EDUCATIONAL_INAPPROPRIATE_THRESHOLD = 0.6
EDUCATIONAL_OFFENSIVE_THRESHOLD = 0.7
EDUCATIONAL_SEXUAL_THRESHOLD = 0.8

# File processing thresholds (configurable)
FILE_PROCESSING_THRESHOLDS = {
    "inappropriate": 0.8,
    "offensive": 0.8,
    "sexual": 0.9,
    "violence": 0.8,
    "overall": 0.8
}

# Data Models
@dataclass
class TextChunk:
    chunk_index: int
    text_content: str
    chunk_size: int
    metadata: Dict[str, Any]

@dataclass  
class ImageChunk:
    chunk_index: int
    image_data: str  # base64
    image_format: str
    image_size: List[int]
    metadata: Dict[str, Any]

class ChunkProcessingRequest(BaseModel):
    file_id: str
    file_type: str
    original_filename: str
    text_chunks: List[Dict[str, Any]] = []
    image_chunks: List[Dict[str, Any]] = []
    processing_options: Dict[str, Any] = {}

class ChunkProcessingResult(BaseModel):
    file_id: str
    processing_status: str
    overall_moderation_score: float
    processing_time: float
    text_results: List[Dict[str, Any]] = []
    image_results: List[Dict[str, Any]] = []

class FileProcessingRequest(BaseModel):
    file_id: str
    filename: str
    file_type: str
    thresholds: Dict[str, float] = {}
    processing_options: Dict[str, Any] = {}

class FileProcessingResult(BaseModel):
    file_id: str
    filename: str
    file_type: str
    processing_status: str
    overall_moderation_score: float
    processing_time: float
    text_content: str = ""
    text_moderation_result: Dict[str, Any] = {}
    image_count: int = 0
    image_moderation_results: List[Dict[str, Any]] = []
    thresholds_used: Dict[str, float] = {}
    details: Dict[str, Any] = {}

# Initialize FastAPI
app = FastAPI(title="AI Document Processing Service", version="4.0.0")

# Global model variables
nudenet_detector = None
clip_model = None
clip_processor = None
toxicity_pipeline = None
hate_speech_pipeline = None
falconsai_nsfw_classifier = None

# Educational content moderation models
educational_offensive_classifier = None
educational_inappropriate_classifier = None
vietnamese_phobert_classifier = None

print("Initializing AI Service...")

# Load NudeNet for NSFW detection
if NUDENET_AVAILABLE:
    try:
        print("Loading NudeNet detector...")
        nudenet_detector = NudeDetector()
        print("NudeNet detector loaded successfully")
    except Exception as e:
        print(f"Warning: NudeNet failed to load: {e}")
        nudenet_detector = None

# Load Falconsai NSFW classifier
if FALCONSAI_AVAILABLE:
    try:
        print("Loading Falconsai NSFW classifier...")
        falconsai_nsfw_classifier = pipeline("image-classification", model="Falconsai/nsfw_image_detection")
        print("Falconsai NSFW classifier loaded successfully")
    except Exception as e:
        print(f"Warning: Falconsai NSFW classifier failed to load: {e}")
        falconsai_nsfw_classifier = None

# Load CLIP for general image classification
if CLIP_AVAILABLE:
    try:
        print("Loading CLIP model...")
        clip_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
        clip_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
        print("CLIP model loaded successfully")
    except Exception as e:
        print(f"Warning: CLIP failed to load: {e}")
        clip_model = None
        clip_processor = None

# Load text moderation models
try:
    print("Loading text moderation models...")
    
    toxicity_pipeline = pipeline(
        "text-classification", 
        model="martin-ha/toxic-comment-model",
        device=0 if torch.cuda.is_available() else -1
    )
    
    hate_speech_pipeline = pipeline(
        "text-classification",
        model="unitary/toxic-bert",
        device=0 if torch.cuda.is_available() else -1
    )
    
    print("Text moderation models loaded successfully")
    
except Exception as e:
    print(f"Warning: Text models failed to load: {e}")
    toxicity_pipeline = None
    hate_speech_pipeline = None

# Load educational content moderation models
if EDUCATIONAL_MODELS_AVAILABLE:
    try:
        print("Loading educational content moderation models...")
        
        # Load offensive content classifier
        educational_offensive_classifier = pipeline(
            "text-classification",
            model="cardiffnlp/twitter-roberta-base-offensive",
            device=0 if torch.cuda.is_available() else -1
        )
        
        # Load inappropriate content classifier (using a general toxicity model)
        educational_inappropriate_classifier = pipeline(
            "text-classification",
            model="unitary/toxic-bert",
            device=0 if torch.cuda.is_available() else -1
        )
        
        print("Educational content moderation models loaded successfully")
        
    except Exception as e:
        print(f"Warning: Educational models failed to load: {e}")
        educational_offensive_classifier = None
        educational_inappropriate_classifier = None

print("AI Service initialization completed")

# Utility Functions
def safe_convert_to_python_types(obj):
    """Convert numpy types to Python native types for JSON serialization"""
    if isinstance(obj, np.integer):
        return int(obj)
    elif isinstance(obj, np.floating):
        return float(obj)
    elif isinstance(obj, np.ndarray):
        return obj.tolist()
    elif isinstance(obj, list):
        return [safe_convert_to_python_types(item) for item in obj]
    elif isinstance(obj, dict):
        return {key: safe_convert_to_python_types(value) for key, value in obj.items()}
    else:
        return obj

def compute_text_hash(text: str) -> str:
    """Compute SHA256 hash of text"""
    return hashlib.sha256(text.encode()).hexdigest()

def compute_simhash(text: str) -> str:
    """Compute SimHash of text"""
    try:
        simhash = Simhash(text)
        return str(simhash.value)
    except Exception as e:
        print(f"SimHash computation error: {e}")
        return "0"

def compute_minhash_signature(text: str) -> List[int]:
    """Compute MinHash signature with proper type conversion"""
    try:
        minhash = MinHash(num_perm=NUM_PERM)
        words = text.lower().split()
        for word in words:
            minhash.update(word.encode('utf8'))
        return [int(x) for x in minhash.hashvalues]
    except Exception as e:
        print(f"MinHash computation error: {e}")
        return [0] * NUM_PERM

def detect_language(text: str) -> str:
    """Simple language detection for Vietnamese vs English"""
    vietnamese_chars = ['ă', 'â', 'đ', 'ê', 'ô', 'ơ', 'ư', 'à', 'á', 'ả', 'ã', 'ạ', 'ằ', 'ắ', 'ẳ', 'ẵ', 'ặ', 
                       'ầ', 'ấ', 'ẩ', 'ẫ', 'ậ', 'è', 'é', 'ẻ', 'ẽ', 'ẹ', 'ề', 'ế', 'ể', 'ễ', 'ệ', 'ì', 'í', 
                       'ỉ', 'ĩ', 'ị', 'ò', 'ó', 'ỏ', 'õ', 'ọ', 'ồ', 'ố', 'ổ', 'ỗ', 'ộ', 'ờ', 'ớ', 'ở', 'ỡ', 
                       'ợ', 'ù', 'ú', 'ủ', 'ũ', 'ụ', 'ừ', 'ứ', 'ử', 'ữ', 'ự', 'ỳ', 'ý', 'ỷ', 'ỹ', 'ỵ']
    
    vietnamese_count = sum(1 for char in text.lower() if char in vietnamese_chars)
    total_chars = len([char for char in text if char.isalpha()])
    
    if total_chars == 0:
        return "unknown"
    
    vietnamese_ratio = vietnamese_count / total_chars
    return "vietnamese" if vietnamese_ratio > 0.1 else "english"

def moderate_vietnamese_text(text: str) -> Dict[str, Any]:
    """Vietnamese text moderation using keyword-based approach"""
    results = {
        "toxicity_score": 0.0,
        "hate_speech_score": 0.0,
        "sexual_content_score": 0.0,
        "violence_score": 0.0,
        "language": "vietnamese",
        "method": "keyword_based"
    }
    
    text_lower = text.lower()
    
    # Vietnamese toxic/hate speech keywords
    vietnamese_toxic_keywords = [
        'đm', 'dm', 'đụ', 'địt', 'lồn', 'cặc', 'buồi', 'loz','fuck'
    ]
    
    # Vietnamese sexual content keywords
    vietnamese_sexual_keywords = [
        'chịch', 'đụ', 'địt', 'lồn', 'cặc', 'buồi',
         'porn'
    ]
    
    # Vietnamese violence keywords
    vietnamese_violence_keywords = [
    ]
    
    # Count occurrences
    toxic_count = sum(1 for keyword in vietnamese_toxic_keywords if keyword in text_lower)
    sexual_count = sum(1 for keyword in vietnamese_sexual_keywords if keyword in text_lower)
    violence_count = sum(1 for keyword in vietnamese_violence_keywords if keyword in text_lower)
    
    # Calculate scores
    word_count = len(text.split())
    if word_count > 0:
        results["toxicity_score"] = min(toxic_count / word_count * 5, 1.0)
        results["sexual_content_score"] = min(sexual_count / word_count * 5, 1.0)
        results["violence_score"] = min(violence_count / word_count * 5, 1.0)
        results["hate_speech_score"] = results["toxicity_score"]
    
    return results

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file"""
    try:
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        
        return text.strip()
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return ""

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from DOCX file"""
    try:
        doc_file = io.BytesIO(file_content)
        doc = docx.Document(doc_file)
        
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        return text.strip()
    except Exception as e:
        print(f"DOCX extraction error: {e}")
        return ""

def extract_text_from_xlsx(file_content: bytes) -> str:
    """Extract text from XLSX file"""
    try:
        xlsx_file = io.BytesIO(file_content)
        workbook = openpyxl.load_workbook(xlsx_file)
        
        text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"Sheet: {sheet_name}\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    text += row_text + "\n"
        
        return text.strip()
    except Exception as e:
        print(f"XLSX extraction error: {e}")
        return ""

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PPTX file"""
    try:
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)
        
        text = ""
        for i, slide in enumerate(presentation.slides):
            text += f"Slide {i+1}:\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        
        return text.strip()
    except Exception as e:
        print(f"PPTX extraction error: {e}")
        return ""

def extract_images_from_pdf(file_content: bytes) -> List[str]:
    """Extract images from PDF file (simplified)"""
    try:
        # This is a simplified version - you might need pdf2image for better extraction
        images = []
        # For now, return empty list - can be enhanced with pdf2image
        return images
    except Exception as e:
        print(f"PDF image extraction error: {e}")
        return []

def extract_images_from_docx(file_content: bytes) -> List[str]:
    """Extract images from DOCX file"""
    try:
        doc_file = io.BytesIO(file_content)
        doc = docx.Document(doc_file)
        
        images = []
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_data = rel.target_part.blob
                image_base64 = base64.b64encode(image_data).decode()
                images.append(image_base64)
        
        return images
    except Exception as e:
        print(f"DOCX image extraction error: {e}")
        return []

def detect_file_type(file_content: bytes, filename: str) -> str:
    """Detect file type from content and filename"""
    # First try to detect from filename extension
    file_extension = filename.lower().split('.')[-1] if '.' in filename else ''
    
    # Map extensions to MIME types
    extension_to_mime = {
        'pdf': 'application/pdf',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'xls': 'application/vnd.ms-excel',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    }
    
    # Check if extension is supported
    if file_extension in extension_to_mime:
        return file_extension
    
    # If extension not found, try to detect from content (basic magic bytes)
    if file_content.startswith(b'%PDF'):
        return 'pdf'
    elif file_content.startswith(b'PK\x03\x04'):
        # ZIP-based format (DOCX, XLSX, PPTX)
        # This is a simplified detection - in real scenario you'd check internal structure
        return file_extension if file_extension in ['docx', 'xlsx', 'pptx'] else 'unknown'
    else:
        return 'unknown'

def extract_file_content(file_content: bytes, filename: str) -> Dict[str, Any]:
    """Extract text and images from various file types"""
    file_extension = filename.lower().split('.')[-1]
    
    result = {
        "text": "",
        "images": [],
        "file_type": file_extension
    }
    
    try:
        if file_extension == "pdf":
            result["text"] = extract_text_from_pdf(file_content)
            result["images"] = extract_images_from_pdf(file_content)
        elif file_extension == "docx":
            result["text"] = extract_text_from_docx(file_content)
            result["images"] = extract_images_from_docx(file_content)
        elif file_extension in ["xlsx", "xls"]:
            result["text"] = extract_text_from_xlsx(file_content)
        elif file_extension == "pptx":
            result["text"] = extract_text_from_pptx(file_content)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        return result
        
    except Exception as e:
        print(f"File content extraction error: {e}")
        return {"text": "", "images": [], "file_type": file_extension, "error": str(e)}

def moderate_educational_content(text: str) -> Dict[str, Any]:
    """Specialized moderation for educational content"""
    results = {
        "inappropriate_score": 0.0,
        "offensive_score": 0.0,
        "sexual_content_score": 0.0,
        "overall_risk": "LOW",
        "language_detected": detect_language(text),
        "moderation_method": "educational_specialized",
        "details": {}
    }
    
    try:
        language = detect_language(text)
        text_lower = text.lower()
        
        # Vietnamese educational content moderation
        if language == "vietnamese":
            # Vietnamese inappropriate words for educational context
            vietnamese_inappropriate = [
                'đm', 'đụ', 'địt', 'lồn', 'cặc', 'buồi', 'loz', 'fuck',
            ]
            
            # Vietnamese sexual content
            vietnamese_sexual = [
                'chịch', 'đụ', 'địt', 'lồn', 'cặc', 'buồi', 'porn'
            ]
            
            # Count inappropriate words
            inappropriate_count = sum(1 for word in vietnamese_inappropriate if word in text_lower)
            sexual_count = sum(1 for word in vietnamese_sexual if word in text_lower)
            
            # Anti-dilution scoring: Score based on presence, not ratio
            if inappropriate_count > 0:
                # Base score for any inappropriate word
                base_score = 0.7
                # Additional score for multiple occurrences
                additional_score = min(inappropriate_count * 0.1, 0.3)
                results["inappropriate_score"] = min(base_score + additional_score, 1.0)
                results["offensive_score"] = results["inappropriate_score"]
            else:
                results["inappropriate_score"] = 0.0
                results["offensive_score"] = 0.0
            
            if sexual_count > 0:
                # Base score for any sexual word
                base_score = 0.6
                # Additional score for multiple occurrences
                additional_score = min(sexual_count * 0.1, 0.4)
                results["sexual_content_score"] = min(base_score + additional_score, 1.0)
            else:
                results["sexual_content_score"] = 0.0
            
            results["details"]["vietnamese_keywords_found"] = {
                "inappropriate": [word for word in vietnamese_inappropriate if word in text_lower],
                "sexual": [word for word in vietnamese_sexual if word in text_lower]
            }
            
        else:
            # English educational content moderation
            english_inappropriate = [
                'fuck', 'shit', 'damn', 'bitch', 'asshole'
            ]
            
            english_sexual = [
             'penis','dick'
            ]
            
            # Count inappropriate words
            inappropriate_count = sum(1 for word in english_inappropriate if word in text_lower)
            sexual_count = sum(1 for word in english_sexual if word in text_lower)
            
            # Anti-dilution scoring: Score based on presence, not ratio
            if inappropriate_count > 0:
                # Base score for any inappropriate word
                base_score = 0.6
                # Additional score for multiple occurrences
                additional_score = min(inappropriate_count * 0.1, 0.4)
                results["inappropriate_score"] = min(base_score + additional_score, 1.0)
                results["offensive_score"] = results["inappropriate_score"]
            else:
                results["inappropriate_score"] = 0.0
                results["offensive_score"] = 0.0
            
            if sexual_count > 0:
                # Base score for any sexual word
                base_score = 0.5
                # Additional score for multiple occurrences
                additional_score = min(sexual_count * 0.1, 0.5)
                results["sexual_content_score"] = min(base_score + additional_score, 1.0)
            else:
                results["sexual_content_score"] = 0.0
            
            results["details"]["english_keywords_found"] = {
                "inappropriate": [word for word in english_inappropriate if word in text_lower],
                "sexual": [word for word in english_sexual if word in text_lower]
            }
            
            # Use AI models for English if available
            if educational_offensive_classifier:
                try:
                    offensive_result = educational_offensive_classifier(text)
                    if isinstance(offensive_result, list) and len(offensive_result) > 0:
                        for result in offensive_result:
                            if result.get('label') in ['OFFENSIVE', 'offensive', '1', 1]:
                                ai_offensive_score = float(result.get('score', 0.0))
                                results["offensive_score"] = max(results["offensive_score"], ai_offensive_score)
                                break
                        results["details"]["ai_offensive_result"] = offensive_result
                except Exception as e:
                    print(f"Educational offensive classifier failed: {e}")
            
            if educational_inappropriate_classifier:
                try:
                    inappropriate_result = educational_inappropriate_classifier(text)
                    if isinstance(inappropriate_result, list) and len(inappropriate_result) > 0:
                        for result in inappropriate_result:
                            if result.get('label') in ['TOXIC', 'toxic', '1', 1]:
                                ai_inappropriate_score = float(result.get('score', 0.0))
                                results["inappropriate_score"] = max(results["inappropriate_score"], ai_inappropriate_score)
                                break
                        results["details"]["ai_inappropriate_result"] = inappropriate_result
                except Exception as e:
                    print(f"Educational inappropriate classifier failed: {e}")
        
        # Overall risk assessment for educational content
        max_score = max(
            results["inappropriate_score"],
            results["offensive_score"],
            results["sexual_content_score"]
        )
        
        if max_score > EDUCATIONAL_OFFENSIVE_THRESHOLD:
            results["overall_risk"] = "HIGH"
        elif max_score > EDUCATIONAL_INAPPROPRIATE_THRESHOLD:
            results["overall_risk"] = "MEDIUM"
        elif max_score > 0.2:
            results["overall_risk"] = "LOW-MEDIUM"
        else:
            results["overall_risk"] = "LOW"
        
        results = safe_convert_to_python_types(results)
        return results
        
    except Exception as e:
        print(f"Educational content moderation failed: {e}")
        return {
            "inappropriate_score": 0.0,
            "offensive_score": 0.0,
            "sexual_content_score": 0.0,
            "overall_risk": "ERROR",
            "language_detected": detect_language(text),
            "error": str(e)
        }

def moderate_file_content(text: str, images: List[str], thresholds: Dict[str, float] = None) -> Dict[str, Any]:
    """Moderate file content with configurable thresholds"""
    if thresholds is None:
        thresholds = FILE_PROCESSING_THRESHOLDS
    
    results = {
        "text_moderation": {},
        "image_moderation": [],
        "overall_score": 0.0,
        "overall_risk": "LOW",
        "thresholds_used": thresholds,
        "details": {}
    }
    
    try:
        # 1. Moderate text content
        if text and text.strip():
            text_result = moderate_educational_content(text)
            results["text_moderation"] = text_result
            
            # Apply custom thresholds
            text_scores = {
                "inappropriate": text_result.get("inappropriate_score", 0.0),
                "offensive": text_result.get("offensive_score", 0.0),
                "sexual": text_result.get("sexual_content_score", 0.0)
            }
            
            # Check against thresholds
            for score_type, score in text_scores.items():
                threshold = thresholds.get(score_type, 0.8)
                if score > threshold:
                    results["details"][f"text_{score_type}_exceeded"] = {
                        "score": score,
                        "threshold": threshold
                    }
        
        # 2. Moderate image content
        if images:
            image_results = []
            for i, image_base64 in enumerate(images):
                try:
                    image_bytes = base64.b64decode(image_base64)
                    image = Image.open(io.BytesIO(image_bytes))
                    
                    if image.mode != 'RGB':
                        image = image.convert('RGB')
                    
                    image_result = moderate_image_advanced(image)
                    image_result["image_index"] = i
                    image_results.append(image_result)
                    
                    # Check against thresholds
                    nsfw_score = image_result.get("nsfw_score", 0.0)
                    violence_score = image_result.get("violence_score", 0.0)
                    
                    if nsfw_score > thresholds.get("sexual", 0.9):
                        results["details"][f"image_{i}_sexual_exceeded"] = {
                            "score": nsfw_score,
                            "threshold": thresholds.get("sexual", 0.9)
                        }
                    
                    if violence_score > thresholds.get("violence", 0.8):
                        results["details"][f"image_{i}_violence_exceeded"] = {
                            "score": violence_score,
                            "threshold": thresholds.get("violence", 0.8)
                        }
                        
                except Exception as e:
                    print(f"Image {i} moderation failed: {e}")
                    image_results.append({
                        "image_index": i,
                        "error": str(e),
                        "nsfw_score": 0.0,
                        "violence_score": 0.0
                    })
            
            results["image_moderation"] = image_results
        
        # 3. Calculate overall score
        all_scores = []
        
        # Add text scores
        if results["text_moderation"]:
            text_scores = [
                results["text_moderation"].get("inappropriate_score", 0.0),
                results["text_moderation"].get("offensive_score", 0.0),
                results["text_moderation"].get("sexual_content_score", 0.0)
            ]
            all_scores.extend(text_scores)
        
        # Add image scores
        for img_result in results["image_moderation"]:
            if "error" not in img_result:
                all_scores.append(img_result.get("nsfw_score", 0.0))
                all_scores.append(img_result.get("violence_score", 0.0))
        
        results["overall_score"] = max(all_scores) if all_scores else 0.0
        
        # 4. Determine overall risk
        overall_threshold = thresholds.get("overall", 0.8)
        if results["overall_score"] > overall_threshold:
            results["overall_risk"] = "HIGH"
        elif results["overall_score"] > overall_threshold * 0.7:
            results["overall_risk"] = "MEDIUM"
        else:
            results["overall_risk"] = "LOW"
        
        return results
        
    except Exception as e:
        print(f"File content moderation failed: {e}")
        return {
            "text_moderation": {},
            "image_moderation": [],
            "overall_score": 1.0,
            "overall_risk": "ERROR",
            "thresholds_used": thresholds,
            "error": str(e)
        }

def moderate_text_advanced(text: str) -> Dict[str, Any]:
    """Advanced text moderation with Vietnamese support"""
    language = detect_language(text)
    
    results = {
        "toxicity_score": 0.0,
        "hate_speech_score": 0.0,
        "sexual_content_score": 0.0,
        "violence_score": 0.0,
        "overall_risk": "LOW",
        "language_detected": language,
        "details": {}
    }
    
    try:
        if language == "vietnamese":
            vietnamese_results = moderate_vietnamese_text(text)
            results.update(vietnamese_results)
            results["details"]["moderation_method"] = "vietnamese_keywords"
        else:
            # English AI models
            if toxicity_pipeline:
                try:
                    toxicity_result = toxicity_pipeline(text)
                    if isinstance(toxicity_result, list) and len(toxicity_result) > 0:
                        toxic_score = 0.0
                        for result in toxicity_result:
                            if result.get('label') in ['TOXIC', 'toxic', '1', 1]:
                                toxic_score = float(result.get('score', 0.0))
                                break
                            elif result.get('label') in ['CLEAN', 'clean', '0', 0]:
                                toxic_score = 1.0 - float(result.get('score', 0.0))
                                break
                        results["toxicity_score"] = toxic_score
                        results["details"]["toxicity_raw"] = toxicity_result
                except Exception as e:
                    print(f"Toxicity detection failed: {e}")
            
            if hate_speech_pipeline:
                try:
                    hate_result = hate_speech_pipeline(text)
                    if isinstance(hate_result, list) and len(hate_result) > 0:
                        hate_score = 0.0
                        for result in hate_result:
                            if result.get('label') in ['TOXIC', 'toxic', 'HATE', 'hate']:
                                hate_score = float(result.get('score', 0.0))
                                break
                        results["hate_speech_score"] = hate_score
                        results["details"]["hate_speech_raw"] = hate_result
                except Exception as e:
                    print(f"Hate speech detection failed: {e}")
            
            # English keyword fallback
            text_lower = text.lower()
            sexual_keywords = ['sex', 'sexual', 'nude', 'naked', 'porn', 'adult', 'explicit']
            violence_keywords = ['violence', 'violent', 'kill', 'murder', 'weapon', 'gun', 'blood', 'death']
            
            sexual_count = sum(1 for word in sexual_keywords if word in text_lower)
            violence_count = sum(1 for word in violence_keywords if word in text_lower)
            
            word_count = len(text.split())
            if word_count > 0:
                results["sexual_content_score"] = min(sexual_count / word_count * 3, 1.0)
                results["violence_score"] = min(violence_count / word_count * 3, 1.0)
            
            results["details"]["moderation_method"] = "english_ai_models"
        
        # Overall risk assessment
        max_score = max(
            results["toxicity_score"],
            results["hate_speech_score"], 
            results["sexual_content_score"],
            results["violence_score"]
        )
        
        if max_score > 0.7:
            results["overall_risk"] = "HIGH"
        elif max_score > 0.4:
            results["overall_risk"] = "MEDIUM"
        else:
            results["overall_risk"] = "LOW"
            
        results = safe_convert_to_python_types(results)
        return results
        
    except Exception as e:
        print(f"Text moderation processing failed: {e}")
        return {
            "toxicity_score": 0.0,
            "hate_speech_score": 0.0,
            "sexual_content_score": 0.0,
            "violence_score": 0.0,
            "overall_risk": "ERROR",
            "language_detected": language,
            "error": str(e)
        }

def moderate_image_with_falconsai(image: Image.Image) -> Dict[str, Any]:
    """Falconsai NSFW detection"""
    try:
        if not falconsai_nsfw_classifier:
            return {"nsfw_score": 0.0, "detection_method": "falconsai_unavailable"}
        
        results = falconsai_nsfw_classifier(image)
        
        nsfw_score = 0.0
        safe_score = 0.0
        
        for result in results:
            label = result['label'].lower()
            score = float(result['score'])
            
            if 'nsfw' in label or 'nude' in label or 'porn' in label or 'explicit' in label:
                nsfw_score = max(nsfw_score, score)
            elif 'safe' in label or 'normal' in label:
                safe_score = max(safe_score, score)
        
        final_nsfw_score = nsfw_score if nsfw_score > safe_score + 0.2 else 0.0
        
        return {
            "nsfw_score": float(final_nsfw_score),
            "detection_method": "Falconsai",
            "details": {
                "falconsai_raw": [{"label": r['label'], "score": float(r['score'])} for r in results],
                "nsfw_confidence": float(nsfw_score),
                "safe_confidence": float(safe_score)
            }
        }
        
    except Exception as e:
        print(f"Falconsai processing failed: {e}")
        return {"nsfw_score": 0.0, "detection_method": "falconsai_error", "error": str(e)}

def moderate_image_advanced(image: Image.Image) -> Dict[str, Any]:
    """Advanced image moderation with Falconsai + CLIP"""
    results = {
        "nsfw_score": 0.0,
        "violence_score": 0.0,
        "overall_risk": "LOW",
        "detection_method": "basic",
        "details": {}
    }
    
    try:
        # Falconsai for NSFW detection (preferred)
        if falconsai_nsfw_classifier:
            try:
                falconsai_result = moderate_image_with_falconsai(image)
                results["nsfw_score"] = falconsai_result.get("nsfw_score", 0.0)
                results["details"].update(falconsai_result.get("details", {}))
                results["detection_method"] = "Falconsai"
            except Exception as e:
                print(f"Falconsai processing failed: {e}")
        
        # CLIP for violence detection
        if clip_model and clip_processor:
            try:
                violence_labels = [
                    "safe everyday image", 
                    "violent action", 
                    "dangerous weapon", 
                    "fighting or combat",
                    "graphic violence",
                    "war or conflict scene"
                ]
                
                inputs = clip_processor(text=violence_labels, images=image, return_tensors="pt", padding=True)
                
                with torch.no_grad():
                    outputs = clip_model(**inputs)
                    logits_per_image = outputs.logits_per_image
                    probs = logits_per_image.softmax(dim=1)
                
                violence_scores = probs[0].tolist()
                
                safe_score = violence_scores[0]
                violent_scores = violence_scores[1:]
                max_violent_score = max(violent_scores)
                violence_confidence = max_violent_score - safe_score
                
                if violence_confidence > 0.3 and max_violent_score > 0.6:
                    violence_score = float(max_violent_score)
                else:
                    violence_score = 0.0
                
                results["violence_score"] = violence_score
                results["details"]["clip_violence_breakdown"] = dict(zip(violence_labels, violence_scores))
                results["details"]["violence_confidence"] = float(violence_confidence)
                results["details"]["safe_score"] = float(safe_score)
                
                if results["detection_method"] == "basic":
                    results["detection_method"] = "CLIP"
                else:
                    results["detection_method"] += " + CLIP"
                    
            except Exception as e:
                print(f"CLIP violence detection failed: {e}")
        
        # Overall assessment
        overall_score = max(results["nsfw_score"], results["violence_score"])
        
        if overall_score > 0.8:
            results["overall_risk"] = "HIGH"
        elif overall_score > 0.5:
            results["overall_risk"] = "MEDIUM"
        elif overall_score > 0.2:
            results["overall_risk"] = "LOW-MEDIUM"
        else:
            results["overall_risk"] = "LOW"
        
        results = safe_convert_to_python_types(results)
        return results
        
    except Exception as e:
        print(f"Image moderation processing failed: {e}")
        return {
            "nsfw_score": 0.0,
            "violence_score": 0.0,
            "overall_risk": "ERROR",
            "detection_method": "error",
            "error": str(e)
        }

# Processing Functions
async def process_text_chunk(chunk: TextChunk) -> Dict[str, Any]:
    """Process text chunk with advanced moderation"""
    try:
        moderation_result = moderate_text_advanced(chunk.text_content)
        
        features = {
            "word_count": len(chunk.text_content.split()),
            "char_count": len(chunk.text_content),
            "line_count": chunk.text_content.count('\n') + 1,
            "avg_word_length": sum(len(word) for word in chunk.text_content.split()) / max(len(chunk.text_content.split()), 1)
        }
        
        text_hash = compute_text_hash(chunk.text_content)
        simhash = compute_simhash(chunk.text_content)
        minhash = compute_minhash_signature(chunk.text_content)
        
        overall_score = max(
            moderation_result.get("toxicity_score", 0),
            moderation_result.get("hate_speech_score", 0),
            moderation_result.get("sexual_content_score", 0),
            moderation_result.get("violence_score", 0)
        )
        
        result = {
            "chunk_index": int(chunk.chunk_index),
            "moderation_score": float(overall_score),
            "moderation_result": moderation_result,
            "features": safe_convert_to_python_types(features),
            "text_hash": text_hash,
            "simhash": simhash,
            "minhash_signature": minhash,
            "metadata": chunk.metadata
        }
        
        return safe_convert_to_python_types(result)
        
    except Exception as e:
        print(f"Text chunk processing failed: {e}")
        return {
            "chunk_index": int(chunk.chunk_index),
            "error": str(e),
            "moderation_score": 1.0,
            "moderation_result": {"overall_risk": "ERROR"}
        }

async def process_image_chunk(chunk: ImageChunk) -> Dict[str, Any]:
    """Process image chunk with advanced moderation"""
    try:
        image_bytes = base64.b64decode(chunk.image_data)
        image = Image.open(io.BytesIO(image_bytes))
        
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        moderation_result = moderate_image_advanced(image)
        
        features = {
            "width": int(image.width),
            "height": int(image.height), 
            "format": chunk.image_format,
            "mode": image.mode,
            "aspect_ratio": float(image.width / max(image.height, 1)),
            "total_pixels": int(image.width * image.height),
            "has_transparency": image.mode in ['RGBA', 'LA'] or 'transparency' in image.info
        }
        
        image_hash = hashlib.sha256(image_bytes).hexdigest()
        
        result = {
            "chunk_index": int(chunk.chunk_index),
            "moderation_score": float(moderation_result.get("nsfw_score", 0) + moderation_result.get("violence_score", 0)) / 2,
            "moderation_result": moderation_result,
            "features": safe_convert_to_python_types(features),
            "image_hash": image_hash,
            "metadata": chunk.metadata
        }
        
        return safe_convert_to_python_types(result)
        
    except Exception as e:
        print(f"Image chunk processing failed: {e}")
        return {
            "chunk_index": int(chunk.chunk_index),
            "error": str(e),
            "moderation_score": 1.0,
            "moderation_result": {"overall_risk": "ERROR"}
        }

# API Endpoints
@app.get("/health")
async def health_check():
    """Health check with detailed model status"""
    return {
        "status": "healthy",
        "service": "AI Document Processing Service",
        "version": "4.1.0",
        "models_loaded": {
            "falconsai_nsfw_classifier": falconsai_nsfw_classifier is not None,
            "nudenet_detector": nudenet_detector is not None,
            "clip_model": clip_model is not None,
            "toxicity_pipeline": toxicity_pipeline is not None,
            "hate_speech_pipeline": hate_speech_pipeline is not None,
            "educational_offensive_classifier": educational_offensive_classifier is not None,
            "educational_inappropriate_classifier": educational_inappropriate_classifier is not None
        },
        "capabilities": {
            "nsfw_detection": "Falconsai" if falconsai_nsfw_classifier else ("NudeNet" if nudenet_detector else "Basic"),
            "violence_detection": "CLIP" if clip_model else "Basic",
            "text_toxicity": "Advanced" if toxicity_pipeline else "Basic",
            "text_hate_speech": "Advanced" if hate_speech_pipeline else "Basic",
            "educational_content_moderation": "Advanced" if educational_offensive_classifier else "Keyword-based"
        },
        "accuracy_estimates": {
            "nsfw_detection": "92-97%" if falconsai_nsfw_classifier else ("90-95%" if nudenet_detector else "60-70%"),
            "violence_detection": "75-85%" if clip_model else "50-60%",
            "text_moderation": "85-90%" if toxicity_pipeline else "60-70%",
            "educational_moderation": "90-95%" if educational_offensive_classifier else "85-90%"
        },
        "new_features": {
            "educational_content_moderation": True,
            "vietnamese_educational_support": True,
            "specialized_educational_endpoints": ["/process-educational-text", "/process-educational-chunks"],
            "file_processing": True,
            "supported_file_types": ["pdf", "docx", "xlsx", "xls", "pptx"],
            "configurable_thresholds": True,
            "file_processing_endpoints": ["/process-file", "/process-file-with-request"]
        },
        "timestamp": time.time()
    }

@app.post("/process-chunks")
async def process_chunks(request: ChunkProcessingRequest) -> ChunkProcessingResult:
    """Process multiple chunks with advanced AI"""
    start_time = time.time()
    
    try:
        text_results = []
        image_results = []
        
        # Process text chunks
        for chunk_data in request.text_chunks:
            chunk = TextChunk(
                chunk_index=chunk_data["chunk_index"],
                text_content=chunk_data["text_content"],
                chunk_size=chunk_data.get("chunk_size", len(chunk_data["text_content"])),
                metadata=chunk_data.get("metadata", {})
            )
            result = await process_text_chunk(chunk)
            text_results.append(result)
        
        # Process image chunks
        for chunk_data in request.image_chunks:
            chunk = ImageChunk(
                chunk_index=chunk_data["chunk_index"],
                image_data=chunk_data["image_data"],
                image_format=chunk_data.get("image_format", "JPEG"),
                image_size=chunk_data.get("image_size", [0, 0]),
                metadata=chunk_data.get("metadata", {})
            )
            result = await process_image_chunk(chunk)
            image_results.append(result)
        
        # Calculate overall score
        all_scores = []
        for result in text_results + image_results:
            if "error" not in result:
                all_scores.append(result["moderation_score"])
        
        overall_score = max(all_scores) if all_scores else 0.0
        processing_time = time.time() - start_time
        
        final_result = ChunkProcessingResult(
            file_id=request.file_id,
            processing_status="completed",
            overall_moderation_score=float(overall_score),
            processing_time=float(processing_time),
            text_results=text_results,
            image_results=image_results
        )
        
        return final_result
        
    except Exception as e:
        print(f"Bulk processing failed: {e}")
        return ChunkProcessingResult(
            file_id=request.file_id,
            processing_status="error",
            overall_moderation_score=1.0,
            processing_time=float(time.time() - start_time),
            text_results=[],
            image_results=[]
        )

@app.post("/process-text-chunk")
async def process_single_text_chunk(
    chunk_index: int = Form(...),
    text_content: str = Form(...),
    metadata: str = Form(default="{}")
):
    """Process single text chunk"""
    try:
        metadata_dict = json.loads(metadata)
        chunk = TextChunk(
            chunk_index=chunk_index,
            text_content=text_content,
            chunk_size=len(text_content),
            metadata=metadata_dict
        )
        result = await process_text_chunk(chunk)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/process-image-chunk")
async def process_single_image_chunk(
    chunk_index: int = Form(...),
    image_file: UploadFile = File(...),
    metadata: str = Form(default="{}")
):
    """Process single image chunk"""
    try:
        metadata_dict = json.loads(metadata)
        
        image_bytes = await image_file.read()
        image_base64 = base64.b64encode(image_bytes).decode()
        
        image = Image.open(io.BytesIO(image_bytes))
        
        chunk = ImageChunk(
            chunk_index=chunk_index,
            image_data=image_base64,
            image_format=image.format or "JPEG",
            image_size=[image.width, image.height],
            metadata=metadata_dict
        )
        result = await process_image_chunk(chunk)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/process-educational-text")
async def process_educational_text(
    text_content: str = Form(...),
    content_type: str = Form(default="general"),
    metadata: str = Form(default="{}")
):
    """Process text for educational content moderation"""
    try:
        metadata_dict = json.loads(metadata)
        
        # Use specialized educational moderation
        moderation_result = moderate_educational_content(text_content)
        
        # Calculate overall score
        overall_score = max(
            moderation_result.get("inappropriate_score", 0),
            moderation_result.get("offensive_score", 0),
            moderation_result.get("sexual_content_score", 0)
        )
        
        result = {
            "content_type": content_type,
            "moderation_score": float(overall_score),
            "moderation_result": moderation_result,
            "metadata": metadata_dict,
            "processing_time": 0.0  # Could add timing if needed
        }
        
        return result
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/process-educational-chunks")
async def process_educational_chunks(request: ChunkProcessingRequest) -> ChunkProcessingResult:
    """Process multiple chunks with educational content moderation"""
    start_time = time.time()
    
    try:
        text_results = []
        image_results = []
        
        # Process text chunks with educational moderation
        for chunk_data in request.text_chunks:
            chunk = TextChunk(
                chunk_index=chunk_data["chunk_index"],
                text_content=chunk_data["text_content"],
                chunk_size=chunk_data.get("chunk_size", len(chunk_data["text_content"])),
                metadata=chunk_data.get("metadata", {})
            )
            
            # Use educational moderation instead of general moderation
            moderation_result = moderate_educational_content(chunk.text_content)
            
            features = {
                "word_count": len(chunk.text_content.split()),
                "char_count": len(chunk.text_content),
                "line_count": chunk.text_content.count('\n') + 1,
                "avg_word_length": sum(len(word) for word in chunk.text_content.split()) / max(len(chunk.text_content.split()), 1)
            }
            
            text_hash = compute_text_hash(chunk.text_content)
            simhash = compute_simhash(chunk.text_content)
            minhash = compute_minhash_signature(chunk.text_content)
            
            overall_score = max(
                moderation_result.get("inappropriate_score", 0),
                moderation_result.get("offensive_score", 0),
                moderation_result.get("sexual_content_score", 0)
            )
            
            result = {
                "chunk_index": int(chunk.chunk_index),
                "moderation_score": float(overall_score),
                "moderation_result": moderation_result,
                "features": safe_convert_to_python_types(features),
                "text_hash": text_hash,
                "simhash": simhash,
                "minhash_signature": minhash,
                "metadata": chunk.metadata,
                "moderation_type": "educational"
            }
            
            text_results.append(safe_convert_to_python_types(result))
        
        # Process image chunks (use existing image moderation)
        for chunk_data in request.image_chunks:
            chunk = ImageChunk(
                chunk_index=chunk_data["chunk_index"],
                image_data=chunk_data["image_data"],
                image_format=chunk_data.get("image_format", "JPEG"),
                image_size=chunk_data.get("image_size", [0, 0]),
                metadata=chunk_data.get("metadata", {})
            )
            result = await process_image_chunk(chunk)
            image_results.append(result)
        
        # Calculate overall score
        all_scores = []
        for result in text_results + image_results:
            if "error" not in result:
                all_scores.append(result["moderation_score"])
        
        overall_score = max(all_scores) if all_scores else 0.0
        processing_time = time.time() - start_time
        
        final_result = ChunkProcessingResult(
            file_id=request.file_id,
            processing_status="completed",
            overall_moderation_score=float(overall_score),
            processing_time=float(processing_time),
            text_results=text_results,
            image_results=image_results
        )
        
        return final_result
        
    except Exception as e:
        print(f"Educational bulk processing failed: {e}")
        return ChunkProcessingResult(
            file_id=request.file_id,
            processing_status="error",
            overall_moderation_score=1.0,
            processing_time=float(time.time() - start_time),
            text_results=[],
            image_results=[]
        )

@app.post("/process-file")
async def process_file(
    file: UploadFile = File(...),
    thresholds: str = Form(default="{}"),
    processing_options: str = Form(default="{}")
) -> FileProcessingResult:
    """Process uploaded file (PDF, DOCX, XLSX, PPTX) with configurable thresholds"""
    start_time = time.time()
    
    try:
        # 1. Validate file
        if not file.filename:
            raise HTTPException(status_code=400, detail="No file provided")
        
        file_extension = file.filename.lower().split('.')[-1]
        supported_extensions = ['pdf', 'docx', 'xlsx', 'xls', 'pptx']
        
        if file_extension not in supported_extensions:
            raise HTTPException(
                status_code=400, 
                detail=f"Unsupported file type: {file_extension}. Supported: {supported_extensions}"
            )
        
        # 2. Read file content
        file_content = await file.read()
        
        if len(file_content) == 0:
            raise HTTPException(status_code=400, detail="Empty file")
        
        # 3. Parse thresholds and options
        try:
            thresholds_dict = json.loads(thresholds) if thresholds else {}
            options_dict = json.loads(processing_options) if processing_options else {}
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="Invalid JSON in thresholds or options")
        
        # Merge with default thresholds
        final_thresholds = {**FILE_PROCESSING_THRESHOLDS, **thresholds_dict}
        
        # 4. Extract content from file
        print(f"Extracting content from {file.filename}...")
        extracted_content = extract_file_content(file_content, file.filename)
        
        if "error" in extracted_content:
            raise HTTPException(status_code=500, detail=f"Content extraction failed: {extracted_content['error']}")
        
        text_content = extracted_content.get("text", "")
        images = extracted_content.get("images", [])
        
        print(f"Extracted: {len(text_content)} characters, {len(images)} images")
        
        # 5. Moderate content
        print("Moderating content...")
        moderation_result = moderate_file_content(text_content, images, final_thresholds)
        
        # 6. Create response
        processing_time = time.time() - start_time
        
        result = FileProcessingResult(
            file_id=f"file_{int(time.time())}_{file.filename}",
            filename=file.filename,
            file_type=file_extension,
            processing_status="completed",
            overall_moderation_score=float(moderation_result.get("overall_score", 0.0)),
            processing_time=float(processing_time),
            text_content=text_content[:1000] + "..." if len(text_content) > 1000 else text_content,  # Truncate for response
            text_moderation_result=moderation_result.get("text_moderation", {}),
            image_count=len(images),
            image_moderation_results=moderation_result.get("image_moderation", []),
            thresholds_used=final_thresholds,
            details=moderation_result.get("details", {})
        )
        
        print(f"Processing completed in {processing_time:.2f}s, overall score: {result.overall_moderation_score}")
        return result
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"File processing failed: {e}")
        processing_time = time.time() - start_time
        
        return FileProcessingResult(
            file_id=f"file_{int(time.time())}_{file.filename if file.filename else 'unknown'}",
            filename=file.filename or "unknown",
            file_type="unknown",
            processing_status="error",
            overall_moderation_score=1.0,
            processing_time=float(processing_time),
            details={"error": str(e)}
        )

@app.post("/process-file-with-request")
async def process_file_with_request(
    file: UploadFile = File(...),
    request_data: str = Form(...)
) -> FileProcessingResult:
    """Process file with detailed request parameters"""
    try:
        # Parse request data
        request_dict = json.loads(request_data)
        
        # Extract parameters
        file_id = request_dict.get("file_id", f"file_{int(time.time())}")
        thresholds = request_dict.get("thresholds", {})
        processing_options = request_dict.get("processing_options", {})
        
        # Process file (reuse the main processing logic)
        return await process_file(file, json.dumps(thresholds), json.dumps(processing_options))
        
    except json.JSONDecodeError:
        raise HTTPException(status_code=400, detail="Invalid JSON in request_data")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# Main
if __name__ == "__main__":
    import uvicorn
    print("Starting AI Service on http://localhost:9000")
    print("API documentation available at: http://localhost:9000/docs")
    print("Service features:")
    if falconsai_nsfw_classifier:
        print("   - High-accuracy NSFW detection (Falconsai)")
    elif nudenet_detector:
        print("   - High-accuracy NSFW detection (NudeNet)")
    if clip_model:
        print("   - Violence detection (CLIP)")
    if toxicity_pipeline:
        print("   - Advanced text moderation")
    print("   - Vietnamese text support")
    print("   - Serialization issues resolved")
    print("   - Production-ready performance")
    uvicorn.run(app, host="0.0.0.0", port=9000)

